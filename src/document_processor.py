import os
from typing import Dict, List, Optional
import PyPDF2
import pandas as pd
from langchain_text_splitters import RecursiveCharacterTextSplitter
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import markdown
from bs4 import BeautifulSoup

class DocumentProcessor:
    """
    Process different types of documents and extract text with multiple fallback methods.
    
    This class handles:
    - PDF processing with multiple extraction methods
    - Text file processing with encoding detection
    - CSV processing with data structure preservation
    - Text chunking for optimal RAG performance
    """
    
    def __init__(self):
        """Initialize document processor with text splitter configuration"""
        # Configure text splitter for optimal chunk sizes
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,      # Maximum characters per chunk
            chunk_overlap=200,    # Overlap to maintain context between chunks
            length_function=len,  # Function to measure chunk length
        )
    
    def process_document(self, file_path: str, filename: str) -> List[Dict]:
        """
        Process a document based on its file extension.
        
        Args:
            file_path: Path to the document file
            filename: Original filename for extension detection
            
        Returns:
            List of chunk dicts: {"text": <chunk text>, "page": <page/slide/
            sheet label or None>}. The "page" label enables page-level
            citations (e.g. "notes.pdf, p.3") when the source format has a
            natural page-like unit (PDF pages, PPTX slides, XLSX sheets);
            it's None for formats without one (TXT, CSV, DOCX, Markdown).
        """
        # Determine processing method based on file extension
        file_extension = filename.lower().split('.')[-1]
        
        if file_extension == 'pdf':
            return self._process_pdf(file_path)
        elif file_extension == 'txt':
            return self._process_txt(file_path)
        elif file_extension == 'csv':
            return self._process_csv(file_path)
        elif file_extension == 'docx':
            return self._process_docx(file_path)
        elif file_extension == 'pptx':
            return self._process_pptx(file_path)
        elif file_extension in ('xlsx', 'xls'):
            return self._process_xlsx(file_path)
        elif file_extension == 'md':
            return self._process_markdown(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
    
    def process_text(self, text: str, source_name: str) -> List[Dict]:
        """
        Process raw text and split into chunks.
        
        Args:
            text: Raw text content
            source_name: Name of the source (for error reporting)
            
        Returns:
            List of chunk dicts: {"text": <chunk text>, "page": None}. No
            page-like unit exists for raw text (e.g. scraped web content),
            so "page" is always None here.
        """
        if not text.strip():
            return []
        
        # Split text into manageable chunks
        return self._wrap_chunks(self.text_splitter.split_text(text), page=None)

    def _wrap_chunks(self, chunks: List[str], page: Optional[object]) -> List[Dict]:
        """
        Tag a list of raw text chunks with a page/slide/sheet label.

        Args:
            chunks: Plain text chunks from the text splitter
            page: Label to attach to every chunk (an int page/slide number,
                  a sheet name string, or None if no such unit applies)

        Returns:
            List of {"text": ..., "page": ...} dicts
        """
        return [{"text": chunk, "page": page} for chunk in chunks]
    
    def _process_pdf(self, file_path: str) -> List[Dict]:
        """
        Extract text from PDF file with multiple fallback methods, tracking
        which page each extracted chunk came from so citations can point
        to a specific page.
        
        This method tries several approaches to handle different PDF types:
        1. PyPDF2 with character cleaning (page-accurate)
        2. pdfplumber (if available, also page-accurate)
        3. Binary extraction fallback (page info lost)
        4. Pattern-based extraction (page info lost)
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            List of {"text": ..., "page": <1-indexed page number or None>} dicts
        """
        # pages_text holds (page_number, cleaned_page_text) for methods that
        # can tell us which page text came from.
        pages_text = []
        
        # Method 1: Try PyPDF2 with character cleaning
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages, start=1):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            # Clean the text to remove problematic characters
                            page_text = self._clean_pdf_text(page_text)
                            if page_text.strip():
                                pages_text.append((page_num, page_text))
                    except Exception as e:
                        print(f"Warning: Could not extract text from page {page_num}: {e}")
                        continue
                        
        except Exception as e:
            print(f"PyPDF2 extraction failed: {e}")
        
        # Method 2: Try pdfplumber if PyPDF2 failed or produced no text
        if not pages_text:
            try:
                import pdfplumber
                print("Trying pdfplumber for PDF extraction...")
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages, start=1):
                        try:
                            page_text = page.extract_text()
                            if page_text:
                                cleaned = self._clean_pdf_text(page_text)
                                if cleaned.strip():
                                    pages_text.append((page_num, cleaned))
                        except Exception as e:
                            print(f"pdfplumber: Could not extract from page {page_num}: {e}")
                            continue
            except ImportError:
                print("pdfplumber not available")
            except Exception as e:
                print(f"pdfplumber extraction failed: {e}")
        
        # If we got page-accurate text, split each page independently so
        # every resulting chunk can be tagged with its correct page number.
        if pages_text:
            chunks = []
            for page_num, page_text in pages_text:
                for chunk in self.text_splitter.split_text(page_text):
                    chunks.append({"text": chunk, "page": page_num})
            if chunks:
                return chunks
        
        # Method 3: Binary extraction fallback (page info is lost here)
        print("Trying binary extraction fallback...")
        text = self._extract_pdf_binary_fallback(file_path)
        
        # Method 4: OCR-like text pattern extraction (page info is lost here)
        if not text.strip():
            print("Trying pattern-based text extraction...")
            text = self._extract_pdf_pattern_fallback(file_path)
        
        # Final validation
        if not text.strip():
            raise ValueError("No text could be extracted from the PDF. The PDF might be image-based, corrupted, or password-protected. Try converting it to a text-based PDF first.")
        
        return self._wrap_chunks(self.text_splitter.split_text(text), page=None)
    
    def _clean_pdf_text(self, text: str) -> str:
        """
        Clean PDF text to handle encoding issues and problematic characters.
        
        Args:
            text: Raw text from PDF
            
        Returns:
            Cleaned text string
        """
        if not text:
            return ""
        
        # Remove problematic Unicode characters
        problematic_chars = [
            '\udbef', '\udcef',  # Surrogate characters
            '\ufeff',            # BOM (Byte Order Mark)
            '\u200b', '\u200c', '\u200d',  # Zero-width characters
            '\u2028', '\u2029',  # Line/paragraph separators
        ]
        
        for char in problematic_chars:
            text = text.replace(char, '')
        
        # Handle encoding issues
        try:
            # Try to encode/decode to clean up encoding issues
            text = text.encode('utf-8', 'ignore').decode('utf-8')
        except:
            # Fallback: keep only printable ASCII characters
            text = ''.join(char for char in text if ord(char) < 128 and (char.isprintable() or char.isspace()))
        
        # Clean up whitespace
        text = ' '.join(text.split())
        
        return text
    
    def _extract_pdf_binary_fallback(self, file_path: str) -> str:
        """
        Binary extraction fallback for problematic PDFs.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text or empty string
        """
        try:
            with open(file_path, 'rb') as file:
                content = file.read()
            
            # Try to decode as latin-1 first (preserves byte values)
            try:
                content_str = content.decode('latin-1', errors='ignore')
            except:
                content_str = content.decode('utf-8', errors='ignore')
            
            # Look for text between common PDF text markers
            import re
            
            # Find text streams in PDF
            text_patterns = []
            
            # Pattern 1: Text between BT and ET markers
            bt_et_pattern = r'BT\s+(.*?)\s+ET'
            matches = re.findall(bt_et_pattern, content_str, re.DOTALL)
            for match in matches:
                # Extract text from PDF text commands
                text_commands = re.findall(r'\((.*?)\)', match)
                text_patterns.extend(text_commands)
            
            # Pattern 2: Look for readable text sequences
            readable_text = re.findall(r'[A-Za-z0-9\s\.,;:!?\-]{20,}', content_str)
            text_patterns.extend(readable_text)
            
            # Combine and clean
            extracted_text = ' '.join(text_patterns)
            return self._clean_pdf_text(extracted_text) if extracted_text else ""
            
        except Exception as e:
            print(f"Binary extraction failed: {e}")
            return ""
    
    def _extract_pdf_pattern_fallback(self, file_path: str) -> str:
        """
        Pattern-based extraction for difficult PDFs.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text or empty string
        """
        try:
            with open(file_path, 'rb') as file:
                content = file.read()
            
            # Convert to string with error handling
            content_str = content.decode('latin-1', errors='ignore')
            
            import re
            
            # Look for various text patterns in PDF structure
            patterns = [
                r'/Title\s*\((.*?)\)',  # PDF title
                r'/Subject\s*\((.*?)\)',  # PDF subject
                r'/Author\s*\((.*?)\)',  # PDF author
                r'>\s*([A-Za-z][A-Za-z0-9\s\.,;:!?\-]{10,})\s*<',  # Text between angle brackets
                r'\]\s*([A-Za-z][A-Za-z0-9\s\.,;:!?\-]{10,})\s*\[',  # Text between square brackets
            ]
            
            extracted_texts = []
            for pattern in patterns:
                matches = re.findall(pattern, content_str, re.IGNORECASE)
                extracted_texts.extend(matches)
            
            # Also try to find any readable text sequences
            readable_sequences = re.findall(r'[A-Za-z][A-Za-z0-9\s\.,;:!?\-]{15,}', content_str)
            extracted_texts.extend(readable_sequences[:10])  # Limit to avoid noise
            
            combined_text = ' '.join(extracted_texts)
            return self._clean_pdf_text(combined_text) if combined_text else ""
            
        except Exception as e:
            print(f"Pattern extraction failed: {e}")
            return ""
    
    def _process_txt(self, file_path: str) -> List[Dict]:
        """
        Process plain text file with multiple encoding support.
        
        Args:
            file_path: Path to text file
            
        Returns:
            List of {"text": ..., "page": None} chunk dicts
        """
        try:
            # Try multiple encodings to handle different file types
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'ascii']
            text = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            # Last resort: read as binary and decode with error handling
            if text is None:
                with open(file_path, 'rb') as file:
                    raw_content = file.read()
                    text = raw_content.decode('utf-8', errors='ignore')
            
            if not text.strip():
                raise ValueError("The text file is empty")
            
            return self._wrap_chunks(self.text_splitter.split_text(text), page=None)
            
        except Exception as e:
            raise ValueError(f"Error processing text file: {str(e)}")
    
    def _process_csv(self, file_path: str) -> List[Dict]:
        """
        Process CSV file with encoding detection and structure preservation.
        
        Args:
            file_path: Path to CSV file
            
        Returns:
            List of {"text": ..., "page": None} chunk dicts representing CSV data
        """
        try:
            # Try multiple encodings for CSV files
            encodings = ['utf-8', 'latin-1', 'cp1252', 'utf-16']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    break
                except (UnicodeDecodeError, UnicodeError):
                    continue
                except Exception as e:
                    if "encoding" not in str(e).lower():
                        raise
            
            if df is None:
                raise ValueError("Could not read CSV file with any supported encoding")
            
            if df.empty:
                raise ValueError("The CSV file is empty")
            
            # Convert DataFrame to text representation
            text_parts = []
            
            # Add column headers and basic info
            text_parts.append("CSV Data Structure:")
            text_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
            text_parts.append(f"Total rows: {len(df)}")
            text_parts.append("\nData:")
            
            # Convert each row to text (limit to first 100 rows for performance)
            max_rows = min(100, len(df))
            for index, row in df.head(max_rows).iterrows():
                row_text = f"Row {index + 1}: "
                row_items = []
                for col, value in row.items():
                    # Clean the value to avoid encoding issues
                    if pd.notna(value):
                        clean_value = str(value).encode('ascii', 'ignore').decode('ascii')
                    else:
                        clean_value = 'N/A'
                    row_items.append(f"{col}: {clean_value}")
                row_text += ", ".join(row_items)
                text_parts.append(row_text)
            
            # Add note if there are more rows
            if len(df) > max_rows:
                text_parts.append(f"\n... and {len(df) - max_rows} more rows")
            
            # Join all parts into final text
            full_text = "\n".join(text_parts)
            
            return self._wrap_chunks(self.text_splitter.split_text(full_text), page=None)
            
        except Exception as e:
            raise ValueError(f"Error processing CSV file: {str(e)}")
    
    def _process_docx(self, file_path: str) -> List[Dict]:
        """
        Extract text from a Word (.docx) document, including paragraphs and tables.
        
        Args:
            file_path: Path to DOCX file
            
        Returns:
            List of {"text": ..., "page": None} chunk dicts. DOCX doesn't
            expose a reliable page count (pagination depends on the reader/
            renderer), so no page label is attached.
        """
        try:
            doc = DocxDocument(file_path)
            text_parts = []
            
            # Extract paragraph text
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract table content (rows joined with " | ")
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            full_text = "\n".join(text_parts)
            
            if not full_text.strip():
                raise ValueError("The DOCX file appears to be empty")
            
            return self._wrap_chunks(self.text_splitter.split_text(full_text), page=None)
            
        except Exception as e:
            raise ValueError(f"Error processing DOCX file: {str(e)}")
    
    def _process_pptx(self, file_path: str) -> List[Dict]:
        """
        Extract text from a PowerPoint (.pptx) presentation, slide by slide,
        tagging each resulting chunk with its slide number for citations.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            List of {"text": ..., "page": <slide number>} chunk dicts
        """
        try:
            prs = Presentation(file_path)
            slides_text = []  # (slide_number, slide_text)
            
            for i, slide in enumerate(prs.slides, start=1):
                slide_lines = []
                
                for shape in slide.shapes:
                    # Extract text from text frames (titles, bullet points, etc.)
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            run_text = "".join(run.text for run in para.runs)
                            if run_text.strip():
                                slide_lines.append(run_text)
                    
                    # Extract text from tables on the slide
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text for cell in row.cells)
                            if row_text.strip():
                                slide_lines.append(row_text)
                
                # Only keep slides that had actual content
                if slide_lines:
                    slides_text.append((i, "\n".join(slide_lines)))
            
            if not slides_text:
                raise ValueError("No text could be extracted from the PPTX. It might contain only images.")
            
            chunks = []
            for slide_num, slide_text in slides_text:
                for chunk in self.text_splitter.split_text(slide_text):
                    chunks.append({"text": chunk, "page": slide_num})
            
            return chunks
            
        except Exception as e:
            raise ValueError(f"Error processing PPTX file: {str(e)}")
    
    def _process_xlsx(self, file_path: str) -> List[Dict]:
        """
        Extract data from an Excel (.xlsx) workbook, sheet by sheet, tagging
        each resulting chunk with its sheet name for citations.
        
        Args:
            file_path: Path to XLSX file
            
        Returns:
            List of {"text": ..., "page": <sheet name>} chunk dicts
        """
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            sheets_text = []  # (sheet_name, sheet_text)
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts = []
                
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    row_vals = [str(cell) for cell in row if cell is not None]
                    if row_vals:
                        text_parts.append(" | ".join(row_vals))
                        row_count += 1
                    # Limit rows per sheet for performance, similar to CSV handling
                    if row_count >= 500:
                        text_parts.append(f"... (sheet truncated after {row_count} rows)")
                        break
                
                if text_parts:
                    sheets_text.append((sheet_name, "\n".join(text_parts)))
            
            if not sheets_text:
                raise ValueError("The XLSX file appears to be empty")
            
            chunks = []
            for sheet_name, sheet_text in sheets_text:
                for chunk in self.text_splitter.split_text(sheet_text):
                    chunks.append({"text": chunk, "page": sheet_name})
            
            return chunks
            
        except Exception as e:
            raise ValueError(f"Error processing XLSX file: {str(e)}")
    
    def _process_markdown(self, file_path: str) -> List[Dict]:
        """
        Extract plain text from a Markdown (.md) file by rendering to HTML
        and stripping tags, so formatting syntax doesn't pollute the context.
        
        Args:
            file_path: Path to Markdown file
            
        Returns:
            List of {"text": ..., "page": None} chunk dicts
        """
        try:
            # Try multiple encodings, consistent with _process_txt
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            raw_text = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        raw_text = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if raw_text is None:
                with open(file_path, 'rb') as f:
                    raw_text = f.read().decode('utf-8', errors='ignore')
            
            # Convert markdown to HTML, then strip tags to get clean text
            html = markdown.markdown(raw_text, extensions=['tables', 'fenced_code'])
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text(separator='\n')
            
            if not text.strip():
                raise ValueError("The Markdown file is empty")
            
            return self._wrap_chunks(self.text_splitter.split_text(text), page=None)
            
        except Exception as e:
            raise ValueError(f"Error processing Markdown file: {str(e)}")
    
    def get_file_info(self, file_path: str) -> dict:
        """
        Get basic information about a file.
        
        Args:
            file_path: Path to file
            
        Returns:
            Dictionary with file information
        """
        try:
            file_size = os.path.getsize(file_path)
            file_name = os.path.basename(file_path)
            file_extension = file_name.split('.')[-1].lower()
            
            return {
                "name": file_name,
                "size": file_size,
                "extension": file_extension,
                "size_mb": round(file_size / (1024 * 1024), 2)
            }
        except Exception as e:
            return {"error": str(e)}